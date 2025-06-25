import os
import re
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
from google import genai
import pytesseract
import random
import io

# Configure Tesseract path based on environment
def configure_tesseract():
    """Configure Tesseract OCR path based on the operating system"""
    if platform.system() == "Windows":
        # Windows path (for local development)
        tesseract_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        if os.path.exists(tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        else:
            print("Warning: Tesseract not found at expected Windows path")
    elif platform.system() == "Linux":
        # Linux path (for Docker container)
        tesseract_path = shutil.which('tesseract')
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        else:
            print("Warning: Tesseract not found in PATH")
    else:
        # macOS or other systems
        tesseract_path = shutil.which('tesseract')
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path

# Initialize Tesseract configuration
configure_tesseract()

def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()

def extract_text_pdf_with_ocr(file_path, word_limit=200):
    try:
        doc = fitz.open(file_path)
        word_list = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if not text.strip():
                # OCR fallback
                pix = page.get_pixmap(dpi=300)
                image_bytes = pix.tobytes("png")
                image = Image.open(io.BytesIO(image_bytes))
                text = pytesseract.image_to_string(image)

            text = clean_text(text)
            words = text.split()
            if not words:
                continue
            word_list.extend(words)
            if len(word_list) >= word_limit:
                break
        return ' '.join(word_list[:word_limit])
    except Exception as e:
        print(f"Error in extract_text_pdf_with_ocr: {e}")
        return ""

def extract_text_pdf_random(file_path, word_limit=1000):
    try:
        doc = fitz.open(file_path)
        words = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if not text.strip():
                break
            text = clean_text(text)
            words += text.split()
            if len(words) >= word_limit:
                break

        if words:
            return ' '.join(words[:word_limit])
        else:
            l = len(doc)
            if l >= 2:
                pages = [0, 1]
            else:
                pages = [0]
            while len(set(pages)) < 4:
                pages.append(random.randint(0, l - 1))
            for i in pages:
                page = doc[i]
                try:
                    pix = page.get_pixmap(dpi=300)
                    image_bytes = pix.tobytes("png")
                    image = Image.open(io.BytesIO(image_bytes))
                    text = pytesseract.image_to_string(image)
                except Exception as e:
                    print(f"OCR failed on page {i}: {e}")
                    continue

                if not text.strip():
                    continue
                text = clean_text(text)
                words += text.split()
                if len(words) >= word_limit:
                    break
        return ' '.join(words[:word_limit]) if words else "No text found in the PDF."
    except Exception as e:
        print(f"Error in extract_text_pdf_random: {e}")
        return ""

def extract_text_docx(file_path, word_limit=200):
    try:
        doc = Document(file_path)
        words = []
        for para in doc.paragraphs:
            text = clean_text(para.text)
            words += text.split()
            if len(words) >= word_limit:
                break
        return ' '.join(words[:word_limit])
    except Exception as e:
        print(f"Error in extract_text_docx: {e}")
        return ""

def extract_text_pptx(file_path, word_limit=200):
    try:
        prs = Presentation(file_path)
        words = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = clean_text(shape.text)
                    words += text.split()
                    if len(words) >= word_limit:
                        break
        return ' '.join(words[:word_limit])
    except Exception as e:
        print(f"Error in extract_text_pptx: {e}")
        return ""

def extract_text_txt(file_path, word_limit=200):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = clean_text(f.read())
            return ' '.join(text.split()[:word_limit])
    except Exception as e:
        print(f"Error in extract_text_txt: {e}")
        return ""

def extract_text_from_file(file_path, word_limit=1000):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.pdf':
            return extract_text_pdf_random(file_path, word_limit)
        elif ext == '.docx':
            return extract_text_docx(file_path, word_limit)
        elif ext == '.pptx':
            return extract_text_pptx(file_path, word_limit)
        elif ext == '.txt':
            return extract_text_txt(file_path, word_limit)
        else:
            raise ValueError("Unsupported file type: " + ext)
    except Exception as e:
        print(f"Error in extract_text_from_file: {e}")
        return ""

def generate_description_from_text(text):
    if not text or len(text.split()) < 15:
        return "No sufficient content available to generate a description."

    try:
        client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
        prompt = (
            "I will provide you with a small portion of text from some study material. "
            "Your task is to analyze the content and generate a short, clear description of what the full material is likely about. "
            "The description should be 2–3 sentences long, summarizing the overall topic and purpose, and suitable as a preview or caption for a notes-sharing platform. "
            "Use simple, professional language that helps students quickly understand what the content covers. "
            "Do not copy exact sentences—paraphrase instead. Avoid unnecessary details and stay focused on the main subject. "

            "If the input is empty or too short to understand the context, return this message: 'Not enough data to generate a description.' "

            "If the content seems to be from a question bank, return a description like: 'This is a set of question bank for [subject_name], containing important questions for practice and review.' "

            "If the input appears unrelated to academic content, or contains irrelevant or non-syllabus-based material, return the same message: 'Not enough data to generate a description.' "

            "Your only job is to generate a meaningful description or handle the input based on these instructions—do not do anything else."
        )

        response = client.models.generate_content(
            model="gemma-3-27b-it",
            contents=prompt + "\n\n" + text
        )
        return response.text
    except Exception as e:
        print(f"Error while generating summary: {e}")
        return ""

